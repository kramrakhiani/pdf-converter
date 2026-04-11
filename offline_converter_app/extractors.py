from __future__ import annotations

import html
import re
import subprocess
import tempfile
from html.parser import HTMLParser
from pathlib import Path
from typing import Sequence

import fitz
from docx import Document
from pptx import Presentation

from .common import ConversionError, ToolState, clean_text, ensure_file, normalize_ext


class HTMLToTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"p", "div", "br", "li", "ul", "ol", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")
        if tag == "li":
            self.parts.append("- ")

    def handle_data(self, data: str) -> None:
        if data.strip():
            self.parts.append(data)

    def handle_endtag(self, tag: str) -> None:
        if tag in {"p", "div", "li", "ul", "ol", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def get_text(self) -> str:
        text = "".join(self.parts)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


def run_command(args: Sequence[str], *, capture: bool = False) -> subprocess.CompletedProcess[str]:
    try:
        return subprocess.run(list(args), check=True, text=True, capture_output=capture)
    except subprocess.CalledProcessError as exc:
        stderr = exc.stderr.strip() if exc.stderr else ""
        stdout = exc.stdout.strip() if exc.stdout else ""
        detail = stderr or stdout or str(exc)
        raise ConversionError(f"Command failed: {' '.join(args)}\n{detail}") from exc


def read_plain_text(path: Path) -> str:
    return clean_text(path.read_text(encoding="utf-8", errors="ignore"))


def read_html_text(path: Path) -> str:
    parser = HTMLToTextParser()
    parser.feed(path.read_text(encoding="utf-8", errors="ignore"))
    return clean_text(html.unescape(parser.get_text()))


def read_rtf_text(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    raw = re.sub(r"\\'[0-9a-fA-F]{2}", "", raw)
    raw = re.sub(r"\\[a-zA-Z]+\d* ?", "", raw)
    raw = raw.replace("{", "").replace("}", "")
    return clean_text(raw)


def read_docx_text(path: Path) -> str:
    doc = Document(str(path))
    return clean_text("\n".join(paragraph.text for paragraph in doc.paragraphs))


def read_pptx_slides(path: Path) -> list[str]:
    prs = Presentation(str(path))
    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = clean_text(shape.text)
                if text:
                    parts.append(text)
        slides.append("\n\n".join(parts).strip() or "(empty slide)")
    return slides


def pdf_text_from_ocr(path: Path, tools: ToolState, lang: str) -> str:
    if not tools.tesseract:
        raise ConversionError("Tesseract is not installed. Install it to OCR scanned PDFs offline.")
    doc = fitz.open(str(path))
    chunks: list[str] = []
    with tempfile.TemporaryDirectory(prefix="offline-converter-") as tmpdir:
        tmpdir_path = Path(tmpdir)
        for index, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=200)
            image_path = tmpdir_path / f"page-{index}.png"
            pix.save(str(image_path))
            result = run_command([tools.tesseract, str(image_path), "stdout", "-l", lang], capture=True)
            text = clean_text(result.stdout)
            if text:
                chunks.append(text)
    return clean_text("\n\n".join(chunks))


def read_pdf_text(path: Path, tools: ToolState, *, ocr: bool = False, ocr_lang: str = "eng") -> str:
    if ocr:
        return pdf_text_from_ocr(path, tools, ocr_lang)
    doc = fitz.open(str(path))
    return clean_text("\n\n".join(page.get_text("text") for page in doc))


def extract_text(path: Path, tools: ToolState, *, ocr: bool = False, ocr_lang: str = "eng") -> str:
    ensure_file(path)
    ext = normalize_ext(path)
    if ext in {"txt", "md", "markdown"}:
        return read_plain_text(path)
    if ext in {"html", "htm"}:
        return read_html_text(path)
    if ext == "rtf":
        return read_rtf_text(path)
    if ext == "docx":
        return read_docx_text(path)
    if ext == "pdf":
        return read_pdf_text(path, tools, ocr=ocr, ocr_lang=ocr_lang)
    if ext == "pptx":
        return clean_text("\n\n".join(read_pptx_slides(path)))
    raise ConversionError(f"Text extraction for .{ext} is not supported.")

