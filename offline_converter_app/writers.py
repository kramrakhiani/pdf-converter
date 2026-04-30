from __future__ import annotations

import html
import io
import tempfile
import textwrap
from pathlib import Path
from typing import Iterable

import fitz
from docx import Document
from docx.shared import Pt as DocxPt
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.utils import ImageReader
from reportlab.platypus import Image as RLImage
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from reportlab.pdfgen import canvas

from .common import ConversionError, clean_text
from .extractors import read_pptx_slides

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


def write_txt(text: str, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(clean_text(text) + "\n", encoding="utf-8")


def write_md(text: str, output_path: Path) -> None:
    write_txt(text, output_path)


def write_html(text: str, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    paragraphs = clean_text(text).split("\n\n")
    body = "\n".join(f"<p>{html.escape(block).replace(chr(10), '<br>')}</p>" for block in paragraphs if block)
    output = (
        "<!doctype html>\n<html>\n<head>\n<meta charset=\"utf-8\">\n"
        "<title>Converted Document</title>\n</head>\n<body>\n"
        f"{body}\n</body>\n</html>\n"
    )
    output_path.write_text(output, encoding="utf-8")


def write_rtf(text: str, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    escaped = clean_text(text).replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}")
    escaped = escaped.replace("\n", "\\par\n")
    output_path.write_text("{\\rtf1\\ansi\n" + escaped + "\n}\n", encoding="utf-8")


def write_docx(text: str, output_path: Path) -> None:
    doc = Document()
    for block in clean_text(text).split("\n\n"):
        paragraph = doc.add_paragraph()
        for index, line in enumerate(block.splitlines()):
            if index:
                paragraph.add_run("\n")
            paragraph.add_run(line)
    doc.save(str(output_path))


def write_docx_from_pdf_pages(input_path: Path, output_path: Path, *, dpi: int = 200) -> None:
    pdf_doc = fitz.open(str(input_path))
    try:
        word_doc = Document()
        section = word_doc.sections[0]

        if len(pdf_doc) > 0:
            first_page = pdf_doc[0].rect
            section.page_width = DocxPt(first_page.width)
            section.page_height = DocxPt(first_page.height)

        section.top_margin = DocxPt(18)
        section.bottom_margin = DocxPt(18)
        section.left_margin = DocxPt(18)
        section.right_margin = DocxPt(18)

        with tempfile.TemporaryDirectory(prefix="offline-converter-") as tmpdir:
            tmpdir_path = Path(tmpdir)
            for index, page in enumerate(pdf_doc, start=1):
                pix = page.get_pixmap(dpi=dpi, alpha=False)
                image_path = tmpdir_path / f"page-{index}.png"
                pix.save(str(image_path))
                run = word_doc.add_paragraph().add_run()
                usable_width = section.page_width - section.left_margin - section.right_margin
                run.add_picture(str(image_path), width=usable_width)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        word_doc.save(str(output_path))
    finally:
        pdf_doc.close()


def _docx_paragraph_markup(paragraph) -> str:
    parts: list[str] = []
    for run in paragraph.runs:
        text = html.escape(run.text).replace("\n", "<br/>")
        if not text:
            continue
        if run.bold:
            text = f"<b>{text}</b>"
        if run.italic:
            text = f"<i>{text}</i>"
        if run.underline:
            text = f"<u>{text}</u>"
        parts.append(text)
    return "".join(parts).strip()


def _docx_paragraph_images(paragraph, doc_part, max_width: float) -> list[RLImage]:
    images: list[RLImage] = []
    for run in paragraph.runs:
        embeds = run._r.xpath(".//*[local-name()='blip']/@*[local-name()='embed']")
        for rel_id in embeds:
            part = doc_part.related_parts.get(rel_id)
            if not part:
                continue
            image_bytes = part.blob
            with Image.open(io.BytesIO(image_bytes)) as img:
                width_px, height_px = img.size
            if width_px <= 0 or height_px <= 0:
                continue
            scale = min(1.0, max_width / width_px)
            rl_img = RLImage(io.BytesIO(image_bytes), width=width_px * scale, height=height_px * scale)
            images.append(rl_img)
    return images


def _docx_container_images(container, max_width: float) -> list[RLImage]:
    images: list[RLImage] = []
    part = getattr(container, "part", None)
    if part is None:
        return images
    for paragraph in container.paragraphs:
        images.extend(_docx_paragraph_images(paragraph, part, max_width))
    return images


def _docx_has_explicit_page_break(paragraph) -> bool:
    return bool(paragraph._p.xpath(".//*[local-name()='br' and @*[local-name()='type' and .='page']]"))


def _docx_has_section_break(paragraph) -> bool:
    return bool(paragraph._p.xpath("./*[local-name()='pPr']/*[local-name()='sectPr']"))


def _docx_header_footer_text(container) -> str:
    parts = [clean_text(paragraph.text) for paragraph in container.paragraphs]
    return "\n".join(part for part in parts if part).strip()


def _length_pt(value, default: float) -> float:
    try:
        if value is not None and value.pt is not None:
            return float(value.pt)
    except Exception:
        pass
    return default


def _safe_color(fill_color, default):
    try:
        if fill_color and getattr(fill_color, "type", None) and fill_color.rgb:
            return colors.HexColor(f"#{fill_color.rgb}")
    except Exception:
        pass
    return default


def _ppt_text_color(shape, default=colors.black):
    try:
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            return _safe_color(shape.text_frame.paragraphs[0].runs[0].font.color, default)
    except Exception:
        pass
    return default


def _ppt_font_name(shape) -> str:
    try:
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            run = shape.text_frame.paragraphs[0].runs[0]
            return "Helvetica-Bold" if run.font.bold else "Helvetica"
    except Exception:
        pass
    return "Helvetica"


def _ppt_font_size(shape, default: float = 18) -> float:
    try:
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            run = shape.text_frame.paragraphs[0].runs[0]
            if run.font.size:
                return max(10, min(28, run.font.size.pt))
    except Exception:
        pass
    return default


def _draw_ppt_text(pdf, shape, left: float, top: float, width: float) -> None:
    text_frame = shape.text_frame
    font_size = _ppt_font_size(shape)
    text_obj = pdf.beginText()
    text_obj.setTextOrigin(left + 4, top - 16)
    text_obj.setFont(_ppt_font_name(shape), font_size)
    text_obj.setLeading(font_size + 3)
    for paragraph in text_frame.paragraphs:
        line = clean_text(paragraph.text)
        if not line:
            text_obj.textLine("")
            continue
        wrapped = textwrap.wrap(line, width=max(12, int(width / max(font_size * 0.55, 1))))
        for segment in wrapped or [""]:
            text_obj.textLine(segment)
    pdf.setFillColor(_ppt_text_color(shape))
    pdf.drawText(text_obj)


def _draw_ppt_table(pdf, shape, left: float, top: float) -> None:
    rows = shape.table.rows
    cols = shape.table.columns
    if len(rows) == 0 or len(cols) == 0:
        return
    row_heights = [row.height * 72.0 / 914400.0 for row in rows]
    col_widths = [col.width * 72.0 / 914400.0 for col in cols]
    y = top
    for row_idx, row in enumerate(rows):
        x = left
        current_height = row_heights[row_idx]
        for col_idx, _col in enumerate(cols):
            current_width = col_widths[col_idx]
            cell = shape.table.cell(row_idx, col_idx)
            pdf.setStrokeColor(colors.black)
            pdf.rect(x, y - current_height, current_width, current_height, stroke=1, fill=0)
            text = clean_text(cell.text)
            if text:
                text_obj = pdf.beginText()
                text_obj.setTextOrigin(x + 4, y - 14)
                text_obj.setFont("Helvetica", 10)
                text_obj.setLeading(12)
                wrapped = textwrap.wrap(text, width=max(8, int(current_width / 5.5)))
                for segment in wrapped[: max(1, int(current_height / 12))]:
                    text_obj.textLine(segment)
                pdf.drawText(text_obj)
            x += current_width
        y -= current_height


def _draw_ppt_shape(pdf, shape, slide_h: float, emu_to_pt: float) -> int:
    left = shape.left * emu_to_pt
    top = slide_h - (shape.top * emu_to_pt)
    width = shape.width * emu_to_pt
    height = shape.height * emu_to_pt

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        unsupported = 0
        for subshape in shape.shapes:
            unsupported += _draw_ppt_shape(pdf, subshape, slide_h, emu_to_pt)
        return unsupported

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_stream = io.BytesIO(shape.image.blob)
        pdf.drawImage(
            ImageReader(image_stream),
            left,
            top - height,
            width=width,
            height=height,
            preserveAspectRatio=True,
            mask="auto",
        )
        return 0

    if getattr(shape, "has_table", False):
        _draw_ppt_table(pdf, shape, left, top)
        return 0

    if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX}:
        fill_color = _safe_color(getattr(shape.fill, "fore_color", None), None)
        line_color = _safe_color(getattr(shape.line, "color", None), colors.black)
        if fill_color is not None:
            pdf.setFillColor(fill_color)
            pdf.setStrokeColor(line_color)
            pdf.rect(left, top - height, width, height, stroke=1, fill=1)
        else:
            pdf.setStrokeColor(line_color)
            pdf.rect(left, top - height, width, height, stroke=1, fill=0)

    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        pdf.setStrokeColor(_safe_color(getattr(shape.line, "color", None), colors.black))
        pdf.setLineWidth(1)
        pdf.line(left, top - height, left + width, top)
        return 0

    if getattr(shape, "has_text_frame", False):
        _draw_ppt_text(pdf, shape, left, top, width)
        return 0

    return 1


def write_pdf_from_docx(path: Path, output_path: Path) -> None:
    doc = Document(str(path))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    section0 = doc.sections[0]

    page_width = max(_length_pt(section0.page_width, float(letter[0])), 200.0)
    page_height = max(_length_pt(section0.page_height, float(letter[1])), 200.0)
    left_margin = max(_length_pt(section0.left_margin, 42.0), 24.0)
    right_margin = max(_length_pt(section0.right_margin, 42.0), 24.0)
    top_margin = max(_length_pt(section0.top_margin, 42.0), 24.0)
    bottom_margin = max(_length_pt(section0.bottom_margin, 42.0), 24.0)
    header_distance = _length_pt(section0.header_distance, 18.0)
    footer_distance = _length_pt(section0.footer_distance, 18.0)
    header_text = _docx_header_footer_text(section0.header)
    footer_text = _docx_header_footer_text(section0.footer)
    usable_band_width = page_width - left_margin - right_margin
    header_images = _docx_container_images(section0.header, min(usable_band_width, 140.0))
    footer_images = _docx_container_images(section0.footer, min(usable_band_width, 140.0))
    font_names = sorted(
        {
            run.font.name
            for paragraph in doc.paragraphs
            for run in paragraph.runs
            if run.font.name
        }
    )

    styles = getSampleStyleSheet()
    body_style = ParagraphStyle(
        "DocxBody",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=11,
        leading=15,
        spaceAfter=8,
    )
    heading_styles = {
        "Heading 1": ParagraphStyle("Heading1Docx", parent=styles["Heading1"], fontName="Helvetica-Bold"),
        "Heading 2": ParagraphStyle("Heading2Docx", parent=styles["Heading2"], fontName="Helvetica-Bold"),
        "Heading 3": ParagraphStyle("Heading3Docx", parent=styles["Heading3"], fontName="Helvetica-Bold"),
    }
    bullet_style = ParagraphStyle(
        "DocxBullet",
        parent=body_style,
        leftIndent=16,
        firstLineIndent=-8,
    )

    story: list[object] = []
    paragraphs = iter(doc.paragraphs)
    tables = iter(doc.tables)
    available_width = page_width - left_margin - right_margin

    for child in doc.element.body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]
        if tag == "p":
            paragraph = next(paragraphs)
            markup = _docx_paragraph_markup(paragraph)
            images = _docx_paragraph_images(paragraph, doc.part, available_width)
            paragraph_style = body_style
            if paragraph.paragraph_format.page_break_before:
                story.append(PageBreak())
            if not markup:
                if images:
                    for image in images:
                        story.append(image)
                        story.append(Spacer(1, 8))
                else:
                    story.append(Spacer(1, 8))
                if _docx_has_explicit_page_break(paragraph) or _docx_has_section_break(paragraph):
                    story.append(PageBreak())
                continue
            style_name = paragraph.style.name if paragraph.style else ""
            if style_name in heading_styles:
                paragraph_style = heading_styles[style_name]
                story.append(Paragraph(markup, paragraph_style))
            elif style_name.startswith("List"):
                is_numbered = "Numbered" in style_name or re.match(r"^\d+[\.\)]", markup.strip())
                if is_numbered:
                    match = re.match(r"^(\d+[\.\)])\s*(.*)", markup.strip())
                    if match:
                        number, rest = match.groups()
                        bullet_text = number
                        markup = rest
                    else:
                        bullet_text = None
                else:
                    bullet_text = "•"
                paragraph_style = ParagraphStyle(
                    f"Bullet-{id(paragraph)}",
                    parent=bullet_style,
                )
                story.append(Paragraph(markup, paragraph_style, bulletText=bullet_text))
            else:
                derived_style = ParagraphStyle(
                    f"Para-{id(paragraph)}",
                    parent=body_style,
                )
                fmt = paragraph.paragraph_format
                if fmt.alignment == 1:
                    derived_style.alignment = 1
                elif fmt.alignment == 2:
                    derived_style.alignment = 2
                elif fmt.alignment == 3:
                    derived_style.alignment = 4
                if fmt.left_indent:
                    derived_style.leftIndent = max(0, fmt.left_indent.pt)
                if fmt.right_indent:
                    derived_style.rightIndent = max(0, fmt.right_indent.pt)
                if fmt.first_line_indent:
                    derived_style.firstLineIndent = fmt.first_line_indent.pt
                if fmt.space_before:
                    derived_style.spaceBefore = max(0, fmt.space_before.pt)
                if fmt.space_after:
                    derived_style.spaceAfter = max(0, fmt.space_after.pt)
                if fmt.line_spacing and isinstance(fmt.line_spacing, (int, float)):
                    derived_style.leading = max(derived_style.fontSize + 2, derived_style.fontSize * float(fmt.line_spacing))
                run_sizes = [run.font.size.pt for run in paragraph.runs if run.font.size]
                if run_sizes:
                    derived_style.fontSize = max(9, min(24, max(run_sizes)))
                    derived_style.leading = max(derived_style.leading, derived_style.fontSize + 3)
                paragraph_style = derived_style
                story.append(Paragraph(markup, paragraph_style))
            for image in images:
                story.append(Spacer(1, 4))
                story.append(image)
            story.append(Spacer(1, 4))
            if _docx_has_explicit_page_break(paragraph) or _docx_has_section_break(paragraph):
                story.append(PageBreak())
        elif tag == "tbl":
            table = next(tables)
            data = []
            row_heights = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = clean_text(cell.text) or " "
                    para = Paragraph(cell_text, body_style)
                    row_data.append(para)
                data.append(row_data)
                if row.height and row.height > 0:
                    row_heights.append(row.height / 914400.0 * 72.0)
            if data:
                num_cols = max(len(row) for row in data) if data else 0
                col_widths = [available_width / num_cols] * num_cols if num_cols > 0 else None
                rl_table = Table(data, colWidths=col_widths, rowHeights=row_heights or None)
                table_style = TableStyle(
                    [
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.darkgrey),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e6e6e6")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
                        ("LEFTPADDING", (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                        ("TOPPADDING", (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                        ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                        ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ]
                )
                rl_table.setStyle(table_style)
                rl_table.hAlign = "LEFT"
                story.append(rl_table)
                story.append(Spacer(1, 12))

    pdf = SimpleDocTemplate(
        str(output_path),
        pagesize=(page_width, page_height),
        leftMargin=left_margin,
        rightMargin=right_margin,
        topMargin=max(top_margin, header_distance + 18),
        bottomMargin=max(bottom_margin, footer_distance + 18),
        title=path.stem,
    )

    def draw_header_footer(canv, _doc_template) -> None:
        canv.saveState()
        if header_text:
            text_obj = canv.beginText(left_margin, page_height - header_distance)
            text_obj.setFont("Helvetica", 9)
            text_obj.setLeading(11)
            for line in header_text.splitlines():
                text_obj.textLine(line)
            canv.drawText(text_obj)
        if header_images:
            x = page_width - right_margin
            for image in reversed(header_images[:2]):
                iw = float(image.drawWidth)
                ih = float(image.drawHeight)
                x -= iw
                image.drawOn(canv, x, page_height - header_distance - ih - 4)
                x -= 8
        if footer_text:
            text_obj = canv.beginText(left_margin, max(12, footer_distance))
            text_obj.setFont("Helvetica", 9)
            text_obj.setLeading(11)
            for line in footer_text.splitlines():
                text_obj.textLine(line)
            canv.drawText(text_obj)
        if footer_images:
            x = page_width - right_margin
            for image in reversed(footer_images[:2]):
                iw = float(image.drawWidth)
                ih = float(image.drawHeight)
                x -= iw
                image.drawOn(canv, x, max(12, footer_distance))
                x -= 8
        canv.restoreState()

    pdf.build(
        story or [Paragraph("(empty document)", body_style)],
        onFirstPage=draw_header_footer,
        onLaterPages=draw_header_footer,
    )
    notes: list[str] = []
    if len(doc.sections) > 1:
        notes.append("Multiple DOCX sections detected; fallback PDF uses the first section's page geometry.")
    if font_names:
        preview = ", ".join(font_names[:6])
        extra = "..." if len(font_names) > 6 else ""
        notes.append(f"Original DOCX fonts detected: {preview}{extra}. Fallback PDF uses base PDF fonts.")
    return tuple(notes)


def wrap_lines(text: str, width: int = 95) -> list[str]:
    lines: list[str] = []
    for block in clean_text(text).split("\n"):
        if not block.strip():
            lines.append("")
            continue
        lines.extend(textwrap.wrap(block, width=width) or [""])
    return lines


def write_pdf_pages(page_blocks: Iterable[list[str]], output_path: Path, *, title: str = "Converted Document") -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    pdf = canvas.Canvas(str(output_path), pagesize=letter)
    _, height = letter
    for page_lines in page_blocks:
        text_obj = pdf.beginText(50, height - 50)
        text_obj.setFont("Helvetica", 11)
        text_obj.setLeading(15)
        pdf.setTitle(title)
        for line in page_lines:
            text_obj.textLine(line)
        pdf.drawText(text_obj)
        pdf.showPage()
    pdf.save()


def text_to_pdf_pages(text: str, *, width: int = 95, lines_per_page: int = 45) -> list[list[str]]:
    lines = wrap_lines(text, width=width)
    pages: list[list[str]] = []
    current: list[str] = []
    for line in lines:
        current.append(line)
        if len(current) >= lines_per_page:
            pages.append(current)
            current = []
    if current or not pages:
        pages.append(current)
    return pages


def write_pdf_from_text(text: str, output_path: Path, *, title: str = "Converted Document") -> None:
    write_pdf_pages(text_to_pdf_pages(text), output_path, title=title)


def write_pdf_from_pptx(path: Path, output_path: Path) -> None:
    prs = Presentation(str(path))
    output_path.parent.mkdir(parents=True, exist_ok=True)

    emu_to_pt = 72.0 / 914400.0
    slide_w = prs.slide_width * emu_to_pt
    slide_h = prs.slide_height * emu_to_pt

    pdf = canvas.Canvas(str(output_path), pagesize=(slide_w, slide_h))
    pdf.setTitle(path.stem)

    notes: list[str] = []
    unsupported_shapes = 0

    for slide in prs.slides:
        background_color = colors.white
        try:
            background_color = _safe_color(slide.background.fill.fore_color, colors.white)
        except Exception:
            pass
        pdf.setFillColor(background_color)
        pdf.rect(0, 0, slide_w, slide_h, stroke=0, fill=1)

        for shape in slide.shapes:
            unsupported_shapes += _draw_ppt_shape(pdf, shape, slide_h, emu_to_pt)

        pdf.showPage()

    pdf.save()
    if unsupported_shapes:
        notes.append(f"{unsupported_shapes} PPT shapes could not be rendered exactly in Python fallback mode.")
    return tuple(notes)


def write_image_pdf(input_path: Path, output_path: Path) -> None:
    ext = input_path.suffix.lower().lstrip(".")
    if ext in {"heic", "heif"}:
        try:
            from pillow_heif import open_heif
            image = open_heif(str(input_path))
            if image.mode in ("RGBA", "LA", "P"):
                image = image.convert("RGB")
        except ImportError:
            try:
                import subprocess
                result = subprocess.run(
                    ["magick", str(input_path), str(output_path)],
                    capture_output=True,
                    timeout=30,
                )
                if result.returncode == 0 and output_path.exists():
                    return
            except FileNotFoundError:
                pass
            raise ConversionError(f"HEIC conversion requires pillow-heif. Install with: pip install pillow-heif")
    elif ext == "svg":
        try:
            import subprocess
            result = subprocess.run(
                ["magick", "convert", str(input_path), "-density", "300", str(output_path)],
                capture_output=True,
                timeout=30,
            )
            if result.returncode == 0 and output_path.exists():
                return
            result = subprocess.run(
                ["rsvg-convert", "-d", "300", "-p", "300", "-o", str(output_path), str(input_path)],
                capture_output=True,
                timeout=30,
            )
            if result.returncode == 0 and output_path.exists():
                return
        except FileNotFoundError:
            pass
        try:
            from svglib.svglib import svg2rlg
            from reportlab.graphics import renderPDF
            drawing = svg2rlg(str(input_path))
            renderPDF.drawToFile(drawing, str(output_path))
            return
        except ImportError:
            raise ConversionError(f"SVG conversion requires svglib. Install with: pip install svglib")
    else:
        image = Image.open(str(input_path))
        if image.mode in ("RGBA", "LA", "P"):
            image = image.convert("RGB")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(str(output_path), "PDF", resolution=100.0)


def pdf_to_pptx(input_path: Path, output_path: Path, *, dpi: int = 200) -> None:
    doc = fitz.open(str(input_path))
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        while prs.slides:
            rel_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rel_id)
            del prs.slides._sldIdLst[0]

        blank_layout = prs.slide_layouts[6]
        with tempfile.TemporaryDirectory(prefix="offline-converter-") as tmpdir:
            tmpdir_path = Path(tmpdir)
            for index, page in enumerate(doc, start=1):
                slide = prs.slides.add_slide(blank_layout)
                pix = page.get_pixmap(dpi=dpi, alpha=False)
                image_path = tmpdir_path / f"page-{index}.png"
                pix.save(str(image_path))
                with Image.open(str(image_path)) as img:
                    img_w, img_h = img.size
                slide_w = prs.slide_width
                slide_h = prs.slide_height
                scale = min(slide_w / img_w, slide_h / img_h)
                width = int(img_w * scale)
                height = int(img_h * scale)
                left = int((slide_w - width) / 2)
                top = int((slide_h - height) / 2)
                slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        prs.save(str(output_path))
    finally:
        doc.close()


def text_to_pptx(text: str, output_path: Path, *, title: str = "Converted Deck") -> None:
    prs = Presentation()
    title_layout = prs.slide_layouts[1]
    blocks = [block.strip() for block in clean_text(text).split("\n\n") if block.strip()] or ["(empty document)"]
    chunk_size = 8
    chunks = [blocks[i : i + chunk_size] for i in range(0, len(blocks), chunk_size)]
    for index, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title if index == 1 else f"{title} ({index})"
        body = slide.placeholders[1].text_frame
        body.clear()
        for item in chunk:
            paragraph = body.add_paragraph()
            paragraph.text = item.replace("\n", " ")
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(18)
    prs.save(str(output_path))


def write_pdf_from_xlsx(input_path: Path, output_path: Path) -> tuple[str, ...]:
    if not HAS_OPENPYXL:
        raise ImportError("openpyxl is required for Excel conversion. Install with: pip install openpyxl")
    wb = openpyxl.load_workbook(str(input_path), data_only=True)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    pdf = canvas.Canvas(str(output_path), pagesize=letter)
    page_width, page_height = letter
    left_margin = 40
    top_margin = 40
    notes: list[str] = []
    sheet_count = 0
    for ws in wb.worksheets:
        sheet_count += 1
        if sheet_count > 1:
            pdf.showPage()
        pdf.setFont("Helvetica", 12)
        pdf.drawString(left_margin, page_height - 30, f"Sheet: {ws.title}")
        y = page_height - 60
        max_col_widths = {}
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
            if y < 50:
                pdf.showPage()
                y = page_height - 40
            x = left_margin
            for col_idx, cell_value in enumerate(row):
                if col_idx >= 10:
                    break
                col_width = max_col_widths.get(col_idx, 0)
                cell_str = str(cell_value or "")
                if len(cell_str) > col_width:
                    max_col_widths[col_idx] = min(len(cell_str), 30)
                pdf.setFont("Helvetica", 9)
                pdf.drawString(x, y, cell_str[:30])
                x += max_col_widths.get(col_idx, 80) + 10
            y -= 15
            if row_idx > 100:
                notes.append(f"Sheet '{ws.title}' truncated at 100 rows.")
                break
    pdf.save()
    result_notes = [f"Converted {sheet_count} sheet(s) from Excel."]
    if sheet_count > 3:
        result_notes.append("Consider using native Excel to PDF for complex spreadsheets.")
    return tuple(result_notes)


def write_xlsx_to_txt(input_path: Path, output_path: Path) -> None:
    if not HAS_OPENPYXL:
        raise ImportError("openpyxl is required for Excel conversion. Install with: pip install openpyxl")
    wb = openpyxl.load_workbook(str(input_path), data_only=True)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"=== Sheet: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            row_str = " | ".join(str(cell or "") for cell in row)
            lines.append(row_str)
        lines.append("")
    output_path.write_text("\n".join(lines), encoding="utf-8")


def write_xlsx_to_csv(input_path: Path, output_path: Path) -> None:
    if not HAS_OPENPYXL:
        raise ImportError("openpyxl is required for Excel conversion. Install with: pip install openpyxl")
    wb = openpyxl.load_workbook(str(input_path), data_only=True)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    import csv
    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for ws in wb.worksheets:
            writer.writerow([f"=== Sheet: {ws.title} ==="])
            for row in ws.iter_rows(values_only=True):
                writer.writerow([str(cell or "") for cell in row])


def write_pdf_from_doc(input_path: Path, output_path: Path) -> tuple[str, ...]:
    notes: list[str] = []
    try:
        import subprocess
        result = subprocess.run(
            ["antiword", str(input_path)],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            write_pdf_from_text(result.stdout, output_path, title=input_path.stem)
            notes.append("Used antiword for .doc to text conversion.")
            return tuple(notes)
    except FileNotFoundError:
        pass
    notes.append(".doc fallback: Converting via embedded image for visual fidelity.")
    write_docx_from_pdf_pages_visual(input_path, output_path)
    return tuple(notes)


def write_docx_from_pdf_pages_visual(input_path: Path, output_path: Path) -> None:
    pdf_doc = fitz.open(str(input_path))
    try:
        doc = Document()
        section = doc.sections[0]
        if len(pdf_doc) > 0:
            first_page = pdf_doc[0].rect
            section.page_width = DocxPt(first_page.width)
            section.page_height = DocxPt(first_page.height)
        section.top_margin = DocxPt(18)
        section.bottom_margin = DocxPt(18)
        section.left_margin = DocxPt(18)
        section.right_margin = DocxPt(18)
        with tempfile.TemporaryDirectory(prefix="offline-converter-") as tmpdir:
            tmpdir_path = Path(tmpdir)
            for index, page in enumerate(pdf_doc, start=1):
                pix = page.get_pixmap(dpi=200)
                image_path = tmpdir_path / f"page-{index}.png"
                pix.save(str(image_path))
                run = doc.add_paragraph().add_run()
                usable_width = section.page_width - section.left_margin - section.right_margin
                run.add_picture(str(image_path), width=usable_width)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
    finally:
        pdf_doc.close()
