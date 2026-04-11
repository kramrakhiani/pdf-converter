from __future__ import annotations

from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation

from .common import ConversionError, ConversionReport, normalize_ext, sha256_file


def verify_output(input_path: Path, output_path: Path, report: ConversionReport) -> ConversionReport:
    if not output_path.exists() or output_path.stat().st_size == 0:
        raise ConversionError(f"Output file was not created correctly: {output_path}")

    source_ext = normalize_ext(input_path)
    target_ext = normalize_ext(output_path)
    notes = list(report.notes)

    if target_ext == "pdf":
        doc = fitz.open(str(output_path))
        page_count = len(doc)
        if page_count == 0:
            raise ConversionError("Generated PDF has zero pages.")
        notes.append(f"Verified PDF output: {page_count} page(s).")
        if source_ext == "pptx":
            try:
                prs = Presentation(str(input_path))
                if len(prs.slides) != page_count:
                    notes.append(
                        f"Slide/page count differs: source slides={len(prs.slides)}, output pages={page_count}."
                    )
            except Exception:
                pass
    elif target_ext == "docx":
        doc = Document(str(output_path))
        paragraph_count = len(doc.paragraphs)
        inline_shapes = len(doc.inline_shapes)
        notes.append(
            f"Verified DOCX output: {paragraph_count} paragraph(s), {inline_shapes} inline object(s)."
        )
        if source_ext == "pdf":
            try:
                source_pdf = fitz.open(str(input_path))
                if inline_shapes < len(source_pdf):
                    notes.append(
                        f"Expected at least {len(source_pdf)} page images in the DOCX, found {inline_shapes}."
                    )
            except Exception:
                pass
    elif target_ext == "pptx":
        prs = Presentation(str(output_path))
        slide_count = len(prs.slides)
        if slide_count == 0:
            raise ConversionError("Generated PPTX has zero slides.")
        notes.append(f"Verified PPTX output: {slide_count} slide(s).")
        if source_ext == "pdf":
            try:
                source_pdf = fitz.open(str(input_path))
                if slide_count != len(source_pdf):
                    notes.append(
                        f"Page/slide count differs: source pages={len(source_pdf)}, output slides={slide_count}."
                    )
            except Exception:
                pass
    else:
        notes.append("Basic output existence check passed.")

    return ConversionReport(
        engine=report.engine,
        fidelity=report.fidelity,
        notes=tuple(notes),
        verified=True,
        checksum=sha256_file(output_path),
    )
